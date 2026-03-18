import os
import re
import time
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from datetime import datetime, timedelta

def load_text_data(file_path):
    data = {}
    if not os.path.exists(file_path): return data
    with open(file_path, 'r', encoding='utf-8') as f:
        for line in f:
            if ':' in line:
                key, value = line.split(':', 1)
                k, v = key.strip(), value.strip()
                if v: data[k] = v
    return data

def find_image_file(folder_path, base_name):
    search_name = base_name.replace("#", "").replace(" ", "").lower()
    extensions = ['.png', '.jpg', '.jpeg', '.PNG', '.JPG', '.JPEG']
    try:
        for file_name in os.listdir(folder_path):
            name_only, ext = os.path.splitext(file_name)
            if name_only.replace(" ", "").lower() == search_name and ext.lower() in extensions:
                return os.path.join(folder_path, file_name)
    except: return None
    return None


def replace_text_preserve_style(paragraph, text_data):
    # 1️⃣ paragraph 전체 텍스트 확인
    full_text = "".join(run.text for run in paragraph.runs)

    matches = re.findall(r'#.*?#', full_text)
    if not matches:
        return

    for tag in matches:
        if tag not in text_data:
            continue

        value = text_data[tag]
        if not value:
            continue

        # 2️⃣ run들을 순회하면서 tag 포함된 부분 찾기
        buffer = ""
        start_idx = None
        end_idx = None

        for i, run in enumerate(paragraph.runs):
            buffer += run.text

            if tag in buffer:
                end_idx = i

                # 시작 run 찾기
                temp = ""
                for j in range(i, -1, -1):
                    temp = paragraph.runs[j].text + temp
                    if tag in temp:
                        start_idx = j
                        break
                break

        if start_idx is None or end_idx is None:
            continue

        # 3️⃣ 해당 구간 텍스트 합치기
        combined_text = ""
        for i in range(start_idx, end_idx + 1):
            combined_text += paragraph.runs[i].text

        # 4️⃣ 치환
        new_text = combined_text.replace(tag, value)

        # 5️⃣ 🔥 핵심: 첫 run에만 넣고 나머지 제거
        paragraph.runs[start_idx].text = new_text

        for i in range(start_idx + 1, end_idx + 1):
            paragraph.runs[i].text = ""

def process_shapes(shapes, text_data, folder_path):
    shapes_to_remove = []
    for shape in shapes:
        if shape.has_text_frame:
            clean_content = shape.text_frame.text.replace(" ", "")
            if "#도면" in clean_content:
                match = re.search(r"#\s*도면\s*\d+\s*#", shape.text_frame.text)
                if match:
                    tag = match.group()
                    img_name = tag.replace("#", "").strip()
                    img_path = find_image_file(folder_path, img_name)
                    if img_path:
                        parent = shape._parent
                        target = parent.shapes if hasattr(parent, "shapes") else parent
                        target.add_picture(img_path, shape.left, shape.top, width=shape.width)
                        shapes_to_remove.append(shape)
                        continue
            for paragraph in shape.text_frame.paragraphs:
                replace_text_preserve_style(paragraph, text_data)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        for paragraph in cell.text_frame.paragraphs:
                            replace_text_preserve_style(paragraph, text_data)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            if hasattr(shape, "shapes"):
                process_shapes(shape.shapes, text_data, folder_path)

    for shape in shapes_to_remove:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
        except: pass

def make_report(project_folder, template_path):
    text_data = load_text_data(os.path.join(project_folder, "내용.txt"))
    prs = Presentation(template_path)

    for slide in prs.slides:
        process_shapes(slide.shapes, text_data, project_folder)

    # 9시간 시차 문제 해결: 역산(Offset) 적용
    # 라이브러리가 저장 시 9시간을 더하므로, 여기서 미리 9시간을 뺍니다.
    now = datetime.now()
    adjusted_now = now - timedelta(hours=9) 

    # PPT 내부 메타데이터 강제 업데이트
    prs.core_properties.created = adjusted_now
    prs.core_properties.modified = adjusted_now
    prs.core_properties.last_modified_by = "Admin"

    
    # 파일명 생성 (템플릿 이름 + timestamp)
    base_name = os.path.splitext(os.path.basename(template_path))[0]
    timestamp = now.strftime('%Y%m%d_%H%M%S')
    output_filename = f"{base_name}_{timestamp}.pptx"
    output_path = os.path.join(project_folder, output_filename)
    prs.save(output_path)

    

    # 윈도우 파일 시스템 시간 동기화
    current_ts = time.time()
    os.utime(output_path, (current_ts, current_ts))

    return output_path